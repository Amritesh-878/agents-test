"""Pure text extractors for class materials (.pptx / .pdf / .docx / .txt / .md).

Deterministic file parsing only — no network, no GPU. Each extractor returns
``(source_filename, block_text)`` pairs; blocks that fail the shared quality
gate (empty slides, boilerplate) are dropped before they reach embedding.
"""

from __future__ import annotations

import re
from pathlib import Path
from typing import Callable, Iterable
from zipfile import BadZipFile

from scripts.embed_and_store import is_quality_text

MaterialBlock = tuple[str, str]

_BLANK_LINE_SPLIT = re.compile(r"\n\s*\n")


class MaterialExtractError(ValueError):
    pass


def _finalize_blocks(path: Path, raw_blocks: Iterable[str]) -> list[MaterialBlock]:
    blocks: list[MaterialBlock] = []
    for raw in raw_blocks:
        text = " ".join(raw.split())
        if text and is_quality_text(text):
            blocks.append((path.name, text))
    return blocks


def extract_pptx(path: Path) -> list[MaterialBlock]:
    """One block per slide: all shape text plus the slide's speaker notes."""
    from pptx import Presentation
    from pptx.exc import PackageNotFoundError

    try:
        presentation = Presentation(str(path))
    except (PackageNotFoundError, BadZipFile, KeyError, OSError) as exc:
        raise MaterialExtractError(f"Failed to parse pptx {path.name}: {exc}") from exc

    raw_blocks: list[str] = []
    for slide in presentation.slides:
        parts = [shape.text_frame.text for shape in slide.shapes if shape.has_text_frame]
        if slide.has_notes_slide and slide.notes_slide.notes_text_frame is not None:
            parts.append(slide.notes_slide.notes_text_frame.text)
        raw_blocks.append("\n".join(part for part in parts if part.strip()))
    return _finalize_blocks(path, raw_blocks)


def extract_pdf(path: Path) -> list[MaterialBlock]:
    """One block per page."""
    from pypdf import PdfReader
    from pypdf.errors import PyPdfError

    try:
        reader = PdfReader(str(path))
        raw_blocks = [page.extract_text() or "" for page in reader.pages]
    except (PyPdfError, ValueError, OSError) as exc:
        raise MaterialExtractError(f"Failed to parse pdf {path.name}: {exc}") from exc
    return _finalize_blocks(path, raw_blocks)


def extract_docx(path: Path) -> list[MaterialBlock]:
    """Consecutive non-empty paragraphs form one block (split on blank paragraphs),
    so short bullet lines stay grouped with their section instead of being dropped
    one-by-one by the quality gate."""
    from docx import Document
    from docx.opc.exceptions import PackageNotFoundError

    try:
        document = Document(str(path))
    except (PackageNotFoundError, BadZipFile, KeyError, OSError) as exc:
        raise MaterialExtractError(f"Failed to parse docx {path.name}: {exc}") from exc

    raw_blocks: list[str] = []
    current: list[str] = []
    for paragraph in document.paragraphs:
        if paragraph.text.strip():
            current.append(paragraph.text)
        elif current:
            raw_blocks.append("\n".join(current))
            current = []
    if current:
        raw_blocks.append("\n".join(current))
    return _finalize_blocks(path, raw_blocks)


def extract_txt(path: Path) -> list[MaterialBlock]:
    """Plain text / markdown: blocks are paragraphs separated by blank lines."""
    try:
        text = path.read_text(encoding="utf-8", errors="replace")
    except OSError as exc:
        raise MaterialExtractError(f"Failed to read {path.name}: {exc}") from exc
    return _finalize_blocks(path, _BLANK_LINE_SPLIT.split(text))


_EXTRACTORS: dict[str, Callable[[Path], list[MaterialBlock]]] = {
    ".pptx": extract_pptx,
    ".pdf": extract_pdf,
    ".docx": extract_docx,
    ".txt": extract_txt,
    ".md": extract_txt,
}

SUPPORTED_EXTENSIONS = frozenset(_EXTRACTORS)


def extract_material_file(path: Path) -> list[MaterialBlock]:
    extractor = _EXTRACTORS.get(path.suffix.lower())
    if extractor is None:
        raise MaterialExtractError(
            f"Unsupported material format '{path.suffix}': {path.name}. "
            f"Supported: {', '.join(sorted(SUPPORTED_EXTENSIONS))}"
        )
    if not path.is_file():
        raise MaterialExtractError(f"Material file not found: {path}")
    return extractor(path)


def extract_materials_dir(materials_dir: Path) -> list[MaterialBlock]:
    """Extract every supported file in the folder (non-recursive), sorted by name."""
    if not materials_dir.is_dir():
        raise MaterialExtractError(f"Materials folder not found: {materials_dir}")
    paths = sorted(
        (
            path
            for path in materials_dir.iterdir()
            if path.is_file() and path.suffix.lower() in _EXTRACTORS
        ),
        key=lambda path: path.name.casefold(),
    )
    blocks: list[MaterialBlock] = []
    for path in paths:
        blocks.extend(extract_material_file(path))
    return blocks
