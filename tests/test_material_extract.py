from __future__ import annotations

from pathlib import Path

import pytest

from scripts.utils.material_extract import (
    SUPPORTED_EXTENSIONS,
    MaterialExtractError,
    extract_docx,
    extract_material_file,
    extract_materials_dir,
    extract_pdf,
    extract_pptx,
    extract_txt,
)

GOOD_SLIDE_TEXT = (
    "The supply function links quantity supplied to the price level and its determinants."
)
GOOD_NOTES_TEXT = (
    "Remind students how each determinant shifts the whole supply curve outward."
)
GOOD_PAGE_TEXT = (
    "Determinants of supply include input prices, technology and government policy."
)
GOOD_PARAGRAPH = (
    "Elasticity of supply measures how strongly quantity supplied reacts to price."
)
BOILERPLATE = "Thank You"


# --- fixture builders ---


def build_pptx(path: Path, slide_texts: list[str], notes: str | None = None) -> None:
    from pptx import Presentation
    from pptx.util import Inches

    presentation = Presentation()
    blank_layout = presentation.slide_layouts[6]
    for index, text in enumerate(slide_texts):
        slide = presentation.slides.add_slide(blank_layout)
        if text:
            box = slide.shapes.add_textbox(Inches(1), Inches(1), Inches(8), Inches(1))
            box.text_frame.text = text
        if notes is not None and index == 0:
            slide.notes_slide.notes_text_frame.text = notes
    presentation.save(str(path))


def build_docx(path: Path, paragraphs: list[str]) -> None:
    from docx import Document

    document = Document()
    for text in paragraphs:
        document.add_paragraph(text)
    document.save(str(path))


def build_pdf_bytes(page_texts: list[str]) -> bytes:
    """Hand-assemble a minimal valid PDF (pypdf reads but cannot author PDFs)."""
    header = b"%PDF-1.4\n"
    objects: list[bytes] = []
    page_count = len(page_texts)
    kids = " ".join(f"{4 + 2 * i} 0 R" for i in range(page_count))
    objects.append(b"1 0 obj\n<< /Type /Catalog /Pages 2 0 R >>\nendobj\n")
    objects.append(
        f"2 0 obj\n<< /Type /Pages /Kids [{kids}] /Count {page_count} >>\nendobj\n".encode()
    )
    objects.append(b"3 0 obj\n<< /Type /Font /Subtype /Type1 /BaseFont /Helvetica >>\nendobj\n")
    for i, text in enumerate(page_texts):
        page_num = 4 + 2 * i
        content_num = page_num + 1
        escaped = text.replace("\\", r"\\").replace("(", r"\(").replace(")", r"\)")
        stream = f"BT /F1 12 Tf 72 720 Td ({escaped}) Tj ET".encode()
        objects.append(
            f"{page_num} 0 obj\n<< /Type /Page /Parent 2 0 R /MediaBox [0 0 612 792] "
            f"/Resources << /Font << /F1 3 0 R >> >> /Contents {content_num} 0 R >>\nendobj\n".encode()
        )
        objects.append(
            f"{content_num} 0 obj\n<< /Length {len(stream)} >>\nstream\n".encode()
            + stream
            + b"\nendstream\nendobj\n"
        )
    body = b""
    offsets: list[int] = []
    pos = len(header)
    for obj in objects:
        offsets.append(pos)
        body += obj
        pos += len(obj)
    xref_pos = pos
    count = len(objects) + 1
    xref = f"xref\n0 {count}\n0000000000 65535 f \n".encode()
    for off in offsets:
        xref += f"{off:010d} 00000 n \n".encode()
    trailer = (
        f"trailer\n<< /Size {count} /Root 1 0 R >>\nstartxref\n{xref_pos}\n%%EOF\n".encode()
    )
    return header + body + xref + trailer


def build_pdf(path: Path, page_texts: list[str]) -> None:
    path.write_bytes(build_pdf_bytes(page_texts))


# --- extract_pptx ---


def test_pptx_extracts_slide_text(tmp_path: Path) -> None:
    path = tmp_path / "supply.pptx"
    build_pptx(path, [GOOD_SLIDE_TEXT])
    blocks = extract_pptx(path)
    assert len(blocks) == 1
    assert blocks[0][0] == "supply.pptx"
    assert GOOD_SLIDE_TEXT in blocks[0][1]


def test_pptx_includes_speaker_notes(tmp_path: Path) -> None:
    path = tmp_path / "supply.pptx"
    build_pptx(path, [GOOD_SLIDE_TEXT], notes=GOOD_NOTES_TEXT)
    blocks = extract_pptx(path)
    assert len(blocks) == 1
    assert GOOD_SLIDE_TEXT in blocks[0][1]
    assert GOOD_NOTES_TEXT in blocks[0][1]


def test_pptx_drops_empty_and_boilerplate_slides(tmp_path: Path) -> None:
    path = tmp_path / "deck.pptx"
    build_pptx(path, [GOOD_SLIDE_TEXT, "", BOILERPLATE])
    blocks = extract_pptx(path)
    texts = [text for _, text in blocks]
    assert texts == [GOOD_SLIDE_TEXT]


def test_pptx_corrupt_file_raises(tmp_path: Path) -> None:
    path = tmp_path / "broken.pptx"
    path.write_bytes(b"this is not a powerpoint file")
    with pytest.raises(MaterialExtractError, match="broken.pptx"):
        extract_pptx(path)


# --- extract_pdf ---


def test_pdf_extracts_per_page(tmp_path: Path) -> None:
    path = tmp_path / "notes.pdf"
    build_pdf(path, [GOOD_PAGE_TEXT, GOOD_PARAGRAPH])
    blocks = extract_pdf(path)
    assert [text for _, text in blocks] == [GOOD_PAGE_TEXT, GOOD_PARAGRAPH]
    assert all(source == "notes.pdf" for source, _ in blocks)


def test_pdf_drops_boilerplate_page(tmp_path: Path) -> None:
    path = tmp_path / "notes.pdf"
    build_pdf(path, [GOOD_PAGE_TEXT, BOILERPLATE])
    blocks = extract_pdf(path)
    assert [text for _, text in blocks] == [GOOD_PAGE_TEXT]


def test_pdf_corrupt_file_raises(tmp_path: Path) -> None:
    path = tmp_path / "broken.pdf"
    path.write_bytes(b"not a pdf at all")
    with pytest.raises(MaterialExtractError, match="broken.pdf"):
        extract_pdf(path)


# --- extract_docx ---


def test_docx_groups_paragraphs_into_blocks(tmp_path: Path) -> None:
    path = tmp_path / "module.docx"
    build_docx(
        path,
        [
            "Determinants of supply",
            "Input prices, technology, and expectations shift supply.",
            "",
            GOOD_PARAGRAPH,
        ],
    )
    blocks = extract_docx(path)
    assert len(blocks) == 2
    # Short heading lines survive because they are grouped with their section body.
    assert "Determinants of supply" in blocks[0][1]
    assert "expectations shift supply" in blocks[0][1]
    assert blocks[1][1] == GOOD_PARAGRAPH


def test_docx_drops_boilerplate_only_block(tmp_path: Path) -> None:
    path = tmp_path / "module.docx"
    build_docx(path, [GOOD_PARAGRAPH, "", BOILERPLATE])
    blocks = extract_docx(path)
    assert [text for _, text in blocks] == [GOOD_PARAGRAPH]


def test_docx_corrupt_file_raises(tmp_path: Path) -> None:
    path = tmp_path / "broken.docx"
    path.write_bytes(b"definitely not a word document")
    with pytest.raises(MaterialExtractError, match="broken.docx"):
        extract_docx(path)


# --- extract_txt ---


def test_txt_splits_blocks_on_blank_lines(tmp_path: Path) -> None:
    path = tmp_path / "notes.txt"
    path.write_text(f"{GOOD_PAGE_TEXT}\n\n{GOOD_PARAGRAPH}\n", encoding="utf-8")
    blocks = extract_txt(path)
    assert [text for _, text in blocks] == [GOOD_PAGE_TEXT, GOOD_PARAGRAPH]


def test_txt_drops_empty_and_boilerplate(tmp_path: Path) -> None:
    path = tmp_path / "notes.txt"
    path.write_text(f"{BOILERPLATE}\n\n\n\n{GOOD_PAGE_TEXT}\n", encoding="utf-8")
    blocks = extract_txt(path)
    assert [text for _, text in blocks] == [GOOD_PAGE_TEXT]


def test_txt_empty_file_returns_no_blocks(tmp_path: Path) -> None:
    path = tmp_path / "empty.txt"
    path.write_text("", encoding="utf-8")
    assert extract_txt(path) == []


# --- extract_material_file ---


def test_extract_material_file_dispatches_md(tmp_path: Path) -> None:
    path = tmp_path / "summary.md"
    path.write_text(f"# Supply\n{GOOD_PAGE_TEXT}\n", encoding="utf-8")
    blocks = extract_material_file(path)
    assert len(blocks) == 1
    assert GOOD_PAGE_TEXT in blocks[0][1]


def test_extract_material_file_unsupported_raises(tmp_path: Path) -> None:
    path = tmp_path / "clip.mp4"
    path.write_bytes(b"video")
    with pytest.raises(MaterialExtractError, match="Unsupported material format"):
        extract_material_file(path)


def test_extract_material_file_missing_raises(tmp_path: Path) -> None:
    with pytest.raises(MaterialExtractError, match="not found"):
        extract_material_file(tmp_path / "ghost.txt")


def test_supported_extensions_cover_spec_formats() -> None:
    assert {".pptx", ".pdf", ".docx", ".txt", ".md"} <= set(SUPPORTED_EXTENSIONS)


# --- extract_materials_dir ---


def test_extract_materials_dir_combines_files_sorted(tmp_path: Path) -> None:
    build_pptx(tmp_path / "b_deck.pptx", [GOOD_SLIDE_TEXT])
    (tmp_path / "a_notes.txt").write_text(GOOD_PAGE_TEXT, encoding="utf-8")
    blocks = extract_materials_dir(tmp_path)
    assert [source for source, _ in blocks] == ["a_notes.txt", "b_deck.pptx"]


def test_extract_materials_dir_ignores_unsupported_files(tmp_path: Path) -> None:
    (tmp_path / "notes.txt").write_text(GOOD_PAGE_TEXT, encoding="utf-8")
    (tmp_path / "recording.mp4").write_bytes(b"video")
    blocks = extract_materials_dir(tmp_path)
    assert [source for source, _ in blocks] == ["notes.txt"]


def test_extract_materials_dir_missing_raises(tmp_path: Path) -> None:
    with pytest.raises(MaterialExtractError, match="not found"):
        extract_materials_dir(tmp_path / "no_such_dir")
